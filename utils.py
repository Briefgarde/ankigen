from pptx import Presentation
import os
import regex as re
import tkinter as tk
from tkinter import filedialog, simpledialog, messagebox
import html


def getInfo():
    root = tk.Tk()
    root.withdraw()

    # 3. Ask for a deck name
    deckname = simpledialog.askstring("Input", "Please enter the name of the deck you want to create:")
    if not deckname:
        messagebox.showerror("Missing Name", "You must enter a name.")
        return None

    # Get .pptx location
    pptxloc = filedialog.askopenfilename(
        title="Select a PowerPoint file",
        filetypes=[("PowerPoint files", "*.pptx")]
    )
    if not pptxloc or not pptxloc.endswith('.pptx'):
        messagebox.showerror("Invalid File", "You must select a .pptx file.")
        return None
    
    # Pics location    
    pic_folder = filedialog.askdirectory(
        title="Select the folder containing the"
    )
    if not pic_folder or not os.path.isdir(pic_folder):
        messagebox.showerror("Invalid Folder", "You must select a valid folder.")
        return None
    
    outputfolder = filedialog.askdirectory(
        title="Select an output folder"
    )
    if not outputfolder or not os.path.isdir(outputfolder):
        messagebox.showerror("Invalid Folder", "You must select a valid folder.")
        return None
    
    data = {
        'pptxloc' : pptxloc,
        'picfolder' : pic_folder,
        'deckname' : deckname,
        'outputfolder' : outputfolder
    }
    return data


def is_bullet_paragraph(para):
    if not para.text or para.text.strip() == "":
        return False
    pPr = para._pPr
    if pPr is None:
        return False
    # Any child tag starting with "a:bu" means a bullet or number
    for child in pPr.iterchildren():
        if child.tag.endswith(("buChar", "buAutoNum", "buBlip", "buNone")):
            return child.tag.endswith(("buChar", "buAutoNum", "buBlip"))
    return False


def getNotes(path):
    notes = []

    prs = Presentation(path)
    slides = prs.slides
    better_arrow = html.unescape("&#x2794;")
    replacements = {
        "" : better_arrow
    }

    
    for slide in slides:
        try:
            if slide.has_notes_slide:
                html_parts = []
                current_level = 0

                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    # Normal text (non-bullet)
                    if not is_bullet_paragraph(para):
                        # Close any open lists
                        while current_level > 0:
                            html_parts.append("</ul>")
                            current_level -= 1

                        # Add plain paragraph
                        run_texts = []
                        for run in para.runs:
                            run_text = run.text
                            if run.font.bold:
                                run_text = f'<span style="color: rgb(255, 0, 0);">{run_text}</span>'
                                # run_texts.append(f'<span style="color: rgb(255, 0, 0);">{run.text}</span>')
                            if run.font.italic:
                                run_text = f'<span style="background-color: rgb(0, 255, 0);">{run_text}</span>'
                                # run_texts.append(f'<span style="background-color: rgb(0, 255, 0);">{run.text}</span>')
                            # else:
                            #     run_texts.append(run.text)
                            run_texts.append(run_text)
                        html_parts.append("<p>" + "".join(run_texts) + "</p>")
                        continue

                    # Bullet paragraph
                    level = para.level or 0

                    # Open new lists if level increased
                    while current_level < level + 1:
                        html_parts.append("<ul>")
                        current_level += 1

                    # Close lists if level decreased
                    while current_level > level + 1:
                        html_parts.append("</ul>")
                        current_level -= 1

                    # Add the bullet itself
                    run_texts = []
                    for run in para.runs:
                        run_text = run.text
                        if run.font.bold:
                            run_text = f'<span style="color: rgb(255, 0, 0);">{run_text}</span>'
                            # run_texts.append(f'<span style="color: rgb(255, 0, 0);">{run.text}</span>')
                        if run.font.italic:
                            run_text = f'<span style="background-color: rgb(0, 255, 0);">{run_text}</span>'
                            # run_texts.append(f'<span style="background-color: rgb(0, 255, 0);">{run.text}</span>')
                        # else:
                        #     run_texts.append(run.text)
                        run_texts.append(run_text)
                    html_parts.append("<li>" + "".join(run_texts) + "</li>")

                # Close any open lists at the end
                while current_level > 0:
                    html_parts.append("</ul>")
                    current_level -= 1

                # Apply replacements
                text = "".join(html_parts)
                for k, v in replacements.items():
                    text = text.replace(k, v)

                notes.append(text.strip())
            else:
                notes.append("No notes so far")
        except Exception:
            notes.append("No notes so far")

    return notes


def getPics(picturePath):
    image_extensions = ('.png', '.jpg', '.jpeg')
    pics = [file for file in os.listdir(picturePath) if file.lower().endswith(image_extensions)]
    # Make sure pics are still in order : 
    pics = sorted(pics, key=lambda x: int(re.search(r'\d+', x).group()))

    pics = [os.path.join(picturePath, file) for file in pics]# Make sure the elements are the absolute path to pics

    return pics

def renamePics(pics, prefix):
    new_paths = []

    for path in pics:
        if not os.path.isfile(path):
            print(f"Warning: Skipping '{path}' because it is not a valid file.")
            continue

        dir_name = os.path.dirname(path)
        base_name = os.path.basename(path)
        new_base_name = prefix.replace(" ", "").lower() + base_name
        new_path = os.path.join(dir_name, new_base_name)

        try:
            os.rename(path, new_path)
            new_paths.append(new_path)
        except Exception as e:
            print(f"Error renaming '{path}' to '{new_path}': {e}")

    return new_paths