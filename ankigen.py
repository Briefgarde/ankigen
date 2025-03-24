from pptx import Presentation
import os
import sys
import regex as re
import genanki as anki


def getPics():
    current = os.listdir(os.getcwd())
    try:
        nameppt = [name for name in current if name.endswith(".pptx")][0][:-5]  # select the first folder with same name,                                                                    
    except:
        print("No folder found for the pictures")

    folderPath = os.path.join(os.getcwd(), nameppt)

    image_extensions = ('.png', '.jpg', '.jpeg')
    pics = [file for file in os.listdir(folderPath) if file.lower().endswith(image_extensions)]
    # Make sure pics are still in order : 
    pics = sorted(pics, key=lambda x: int(re.search(r'\d+', x).group()))

    pics = [os.path.join(folderPath, file) for file in pics]# Make sure the elements are the absolute path to pics

    return pics


def getNotes():
    notes = []

    current = os.listdir(os.getcwd())
    nameppt = [name for name in current if name.endswith(".pptx")][0] 
    pathToPrs = os.path.join(os.getcwd(), nameppt)

    prs = Presentation(pathToPrs)
    slides = prs.slides

    for slide in slides:
        if slide.has_notes_slide:
            notes.append(slide.notes_slide.notes_text_frame.text)
        else:
            notes.append("")
    return notes

def genDeck():
    notes = getNotes()
    pics = getPics()
    if len(notes) != len(pics):
        raise Exception("Elements length do not match")
    print("Elements length match")

    for file in os.listdir(os.getcwd()):
        if file.endswith(".apkg"):
            os.remove(file)


    my_model = anki.Model(
        1607392318,
        'Model Leon',
        fields=[
            {'name': 'Count'},
            {'name': 'Note'},
            {'name': 'MyMedia'},                                 
        ],
        templates=[
            {
                'name': 'Card {{Count}}',
                'qfmt': '{{MyMedia}}',             
                'afmt': '{{FrontSide}}<hr id="note">{{Note}}',
                },
            ])
    print("Model generé")

    deck = anki.Deck(
            1777167684, 
            'Cartes Leon'
    )

    # cartes = []
    for i in range(len(notes)):
        picname = os.path.basename(pics[i])
        media_name = f'<img src="{picname}">'
        carte = anki.Note(
                model=my_model,
                fields=[f"{i}", notes[i], media_name] 
        )
        print(f"Carte {i} générée {carte}")
        deck.add_note(carte)
    print("Deck généré")

    package = anki.Package(deck)
    print("Package généré")

    package.media_files = pics
    print("Pics given to medial file")

    package.write_to_file(f"output.apkg")
    print("Deck crée")

if __name__ == '__main__':
    genDeck()