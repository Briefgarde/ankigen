from pptx import Presentation
import os
import sys
import regex as re
import genanki as anki
import tkinter as tk
from tkinter import filedialog, simpledialog, messagebox


def getInfo():
    root = tk.Tk()
    root.withdraw()

    # 3. Ask for a deck name
    deckname = simpledialog.askstring("Input", "Please enter the name of the deck you want to create:")
    if not deckname:
        messagebox.showerror("Missing Name", "You must enter a name.")
        return None

    # Get .pptx location
    pptxloc = filedialog.askopenfilename(
        title="Select a PowerPoint file",
        filetypes=[("PowerPoint files", "*.pptx")]
    )
    if not pptxloc or not pptxloc.endswith('.pptx'):
        messagebox.showerror("Invalid File", "You must select a .pptx file.")
        return None
    
    # Ask is user ok with modifying the names of the files
    consent_text = (
        """
The program will now ask you to locate the folder containing the images from the slides of your powerpoint. If you haven't already, please export all the slide to a picture format (png/jpg) into a single folder
The name of the slides will be modified so that they include the name of the deck you're creating. 
Do you agree to that ? 
        """
    )
    consent = messagebox.askyesno("Consent Required", consent_text)
    if not consent:
        messagebox.showinfo("Consent Denied", "Operation cancelled due to lack of consent.")
        raise Exception("The program can not continue if you do not agree to this.")
    
    # Pics location    
    pic_folder = filedialog.askdirectory(
        title="Select a folder"
    )
    if not pic_folder or not os.path.isdir(pic_folder):
        messagebox.showerror("Invalid Folder", "You must select a valid folder.")
        return None
    
    outputfolder = filedialog.askdirectory(
        title="Select an output folder"
    )
    if not outputfolder or not os.path.isdir(outputfolder):
        messagebox.showerror("Invalid Folder", "You must select a valid folder.")
        return None
    
    data = {
        'pptxloc' : pptxloc,
        'picfolder' : pic_folder,
        'deckname' : deckname,
        'outputfolder' : outputfolder
    }
    return data


def getNotes(path):
    notes = []

    prs = Presentation(path)
    slides = prs.slides

    for slide in slides:
        if slide.has_notes_slide:
            notes.append(slide.notes_slide.notes_text_frame.text)
        else:
            notes.append("No notes so far")
    return notes


def getPics(picturePath):
    image_extensions = ('.png', '.jpg', '.jpeg')
    pics = [file for file in os.listdir(picturePath) if file.lower().endswith(image_extensions)]
    # Make sure pics are still in order : 
    pics = sorted(pics, key=lambda x: int(re.search(r'\d+', x).group()))

    pics = [os.path.join(picturePath, file) for file in pics]# Make sure the elements are the absolute path to pics

    return pics

def renamePics(pics, prefix):
    new_paths = []

    for path in pics:
        if not os.path.isfile(path):
            print(f"Warning: Skipping '{path}' because it is not a valid file.")
            continue

        dir_name = os.path.dirname(path)
        base_name = os.path.basename(path)
        new_base_name = prefix + base_name
        new_path = os.path.join(dir_name, new_base_name)

        try:
            os.rename(path, new_path)
            new_paths.append(new_path)
        except Exception as e:
            print(f"Error renaming '{path}' to '{new_path}': {e}")

    return new_paths

def genDeck(notes, pics, deckname, outputfolder):
    if len(notes) != len(pics):
        raise Exception("Elements length do not match")
    my_model = anki.Model(
        1607392318,
        'Model Leon',
        fields=[
            {'name': 'Count'},
            {'name': 'Note'},
            {'name': 'MyMedia'},                                 
        ],
        templates=[
            {
                'name': 'Card {{Count}}',
                'qfmt': '{{MyMedia}}',             
                'afmt': '{{FrontSide}}<hr id="note">{{Note}}',
                },
            ])

    deck = anki.Deck(
            1777167684, 
            f'{deckname}'
    )
    for i in range(len(notes)):
        picname = os.path.basename(pics[i])
        media_name = f'<img src="{picname}">'
        carte = anki.Note(
                model=my_model,
                fields=[f"{i}", notes[i], media_name] 
        )
        deck.add_note(carte)
    
    package = anki.Package(deck)


    package.media_files = pics

    package.write_to_file(f"{os.path.join(outputfolder, deckname)}.apkg")
    
if __name__ == "__main__":
    data = getInfo()
    if not data:
        raise Exception("Some info were not provided correctly. Please try again")
    notes = getNotes(data['pptxloc'])
    pics = renamePics(getPics(data['picfolder']), data['deckname'])
    genDeck(notes, pics, deckname=data['deckname'], outputfolder=data['outputfolder'])